"""
Local AI File Assistant — Text Extractor
Extracts plain text from PDF, DOCX, PPTX, XLSX, code, and plain-text files.
Each format has a dedicated handler; unsupported types fall back to raw text.
Errors are returned as strings — extractors never raise.
"""

from pathlib import Path

import chardet

import config

_MAX_EXTRACT_CHARS = 500_000


class TextExtractor:
    """Route files to the correct extractor based on suffix."""

    def __init__(self):
        # Build dispatch tables: extension → full-text handler and page handler
        self._handlers: dict[str, callable] = {}
        self._page_handlers: dict[str, callable] = {}

        _map = [
            ("pdf",  self._extract_pdf,  self._pages_pdf),
            ("docx", self._extract_docx, self._pages_docx),
            ("pptx", self._extract_pptx, self._pages_pptx),
            ("xlsx", self._extract_xlsx, self._pages_xlsx),
            ("text", self._extract_text, self._pages_text),
            ("code", self._extract_text, self._pages_text),
        ]
        for group, extractor, pager in _map:
            for ext in config.SUPPORTED_EXTENSIONS[group]:
                self._handlers[ext]      = extractor
                self._page_handlers[ext] = pager

    # ── Public API ───────────────────────────────────────────────────────────

    def extract(self, file_path: Path) -> str:
        """Extract full text from file_path. Returns text or error string — never raises."""
        file_path = Path(file_path)
        if not file_path.is_file():
            return "Error: '{}' is not a file or does not exist.".format(file_path.name)

        handler = self._handlers.get(file_path.suffix.lower(), self._extract_text)
        try:
            text = handler(file_path)
            if len(text) > _MAX_EXTRACT_CHARS:
                text = text[:_MAX_EXTRACT_CHARS] + "\n\n[... truncated ...]"
            return text
        except Exception as exc:
            return "Error extracting '{}': {}".format(file_path.name, exc)

    # Page-aware extraction — returns (page_number, text) pairs for precise operations
    def extract_pages(self, file_path: Path) -> list[tuple[int, str]]:
        """Extract text page by page. Returns [(page_num, text), ...] — never raises."""
        file_path = Path(file_path)
        if not file_path.is_file():
            return [(1, "Error: '{}' not found.".format(file_path.name))]

        handler = self._page_handlers.get(file_path.suffix.lower(), self._pages_text)
        try:
            pages = handler(file_path)
            return pages if pages else [(1, "(No extractable text found)")]
        except Exception as exc:
            return [(1, "Error extracting pages from '{}': {}".format(file_path.name, exc))]

    def supported(self, file_path: Path) -> bool:
        return Path(file_path).suffix.lower() in config.ALL_EXTENSIONS

    # ── Full-text extractors ──────────────────────────────────────────────────

    @staticmethod
    def _extract_pdf(path: Path) -> str:
        import fitz

        pages: list[str] = []
        with fitz.open(path) as doc:
            for i, page in enumerate(doc, start=1):
                text = page.get_text().strip()
                if text:
                    pages.append("--- Page {} ---\n{}".format(i, text))
        return "\n\n".join(pages) if pages else "(No extractable text found)"

    @staticmethod
    def _extract_docx(path: Path) -> str:
        from docx import Document

        doc   = Document(str(path))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        for table in doc.tables:
            rows = ["\t".join(c.text.strip() for c in row.cells) for row in table.rows]
            if rows:
                parts.append("\n".join(rows))
        return "\n\n".join(parts) if parts else "(No extractable text found)"

    @staticmethod
    def _extract_pptx(path: Path) -> str:
        from pptx import Presentation

        slides: list[str] = []
        for i, slide in enumerate(Presentation(str(path)).slides, start=1):
            texts = [
                para.text.strip()
                for shape in slide.shapes if shape.has_text_frame
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            ]
            if texts:
                slides.append("--- Slide {} ---\n{}".format(i, "\n".join(texts)))
        return "\n\n".join(slides) if slides else "(No extractable text found)"

    @staticmethod
    def _extract_xlsx(path: Path) -> str:
        from openpyxl import load_workbook

        wb     = load_workbook(str(path), data_only=True, read_only=True)
        sheets: list[str] = []
        for name in wb.sheetnames:
            rows = [
                "\t".join(str(c) if c is not None else "" for c in row)
                for row in wb[name].iter_rows(values_only=True)
                if any(c is not None for c in row)
            ]
            if rows:
                sheets.append("--- Sheet: {} ---\n{}".format(name, "\n".join(rows)))
        wb.close()
        return "\n\n".join(sheets) if sheets else "(No extractable text found)"

    @staticmethod
    def _extract_text(path: Path) -> str:
        raw      = path.read_bytes()
        sample   = raw[:32_768]
        encoding = chardet.detect(sample).get("encoding") or "utf-8"
        try:
            return raw.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            return raw.decode("utf-8", errors="replace")

    # ── Page-level extractors ─────────────────────────────────────────────────

    @staticmethod
    def _pages_pdf(path: Path) -> list[tuple[int, str]]:
        import fitz

        pages = []
        with fitz.open(path) as doc:
            for i, page in enumerate(doc, start=1):
                text = page.get_text().strip()
                if text:
                    pages.append((i, text))
        return pages

    @staticmethod
    def _pages_docx(path: Path) -> list[tuple[int, str]]:
        from docx import Document

        doc   = Document(str(path))
        items: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                items.append(para.text.strip())
        for table in doc.tables:
            items.append("\n".join("\t".join(c.text.strip() for c in r.cells) for r in table.rows))

        # Group paragraphs into synthetic pages of LINES_PER_PAGE items
        size  = config.LINES_PER_PAGE
        return [
            (i // size + 1, "\n\n".join(items[i : i + size]))
            for i in range(0, len(items), size)
        ]

    @staticmethod
    def _pages_pptx(path: Path) -> list[tuple[int, str]]:
        from pptx import Presentation

        pages = []
        for i, slide in enumerate(Presentation(str(path)).slides, start=1):
            texts = [
                para.text.strip()
                for shape in slide.shapes if shape.has_text_frame
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            ]
            if texts:
                pages.append((i, "\n".join(texts)))
        return pages

    @staticmethod
    def _pages_xlsx(path: Path) -> list[tuple[int, str]]:
        from openpyxl import load_workbook

        wb, pages = load_workbook(str(path), data_only=True, read_only=True), []
        for i, name in enumerate(wb.sheetnames, start=1):
            rows = [
                "\t".join(str(c) if c is not None else "" for c in row)
                for row in wb[name].iter_rows(values_only=True)
                if any(c is not None for c in row)
            ]
            if rows:
                pages.append((i, "\n".join(rows)))
        wb.close()
        return pages

    @staticmethod
    def _pages_text(path: Path) -> list[tuple[int, str]]:
        raw      = path.read_bytes()
        encoding = chardet.detect(raw[:32_768]).get("encoding") or "utf-8"
        try:
            text = raw.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            text = raw.decode("utf-8", errors="replace")

        lines = text.splitlines()
        size  = config.LINES_PER_PAGE
        pages = [
            (i // size + 1, "\n".join(lines[i : i + size]))
            for i in range(0, len(lines), size)
            if any(lines[i : i + size])
        ]
        return pages if pages else [(1, text)]
